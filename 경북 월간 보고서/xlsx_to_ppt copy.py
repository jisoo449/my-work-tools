from __future__ import annotations

from dataclasses import dataclass
from typing import Dict, Optional, Tuple, List

import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

# -----------------------------
# Data model
# -----------------------------
@dataclass
class Stats:
    min: Optional[float] = None
    max: Optional[float] = None
    avg: Optional[float] = None

@dataclass
class MetricBlock:
    # CPU/MEM 은 2개 지표, Network는 IN/OUT 2개 지표 등 복수 가능
    stats_by_metric: Dict[str, Stats]     # e.g. {"CPU": Stats(...), "MEM": Stats(...)}
    images: List[bytes]                  # 이미지(PNG/JPG 바이너리). 블록 내 이미지 순서대로 담김

@dataclass
class ServerReport:
    server_name: str
    cpu_mem: MetricBlock
    network: MetricBlock
    filesystem: MetricBlock

# -----------------------------
# Formatting helpers
# -----------------------------
def _fmt(v: Optional[float], *, decimals: int = 2, suffix: str = "") -> str:
    if v is None:
        return ""
    return f"{v:.{decimals}f}{suffix}"

def _fmt_pct(v: Optional[float]) -> str:
    # 입력이 이미 % 단위(예: 12.34)라고 가정 → "12.34 %"
    return "" if v is None else f"{v:.2f} %"

def _fmt_bps(v: Optional[float]) -> str:
    # bps를 사람이 읽기 좋게 축약 (k, M, G)
    if v is None:
        return ""
    x = float(v)
    units = [("G", 1e9), ("M", 1e6), ("k", 1e3)]
    for u, base in units:
        if abs(x) >= base:
            return f"{x/base:.2f} {u}"
    return f"{x:.2f}"

def _norm_server_name(s: str) -> str:
    # "bastion-lnx (1.2.3.4)" → "bastion-lnx"
    s = (s or "").strip()
    s = re.sub(r"\s*\(.*\)\s*$", "", s)
    return s.strip()


# -----------------------------
# PPT shape search helpers
# -----------------------------
def find_shape_by_name(slide, name: str):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None

def find_text_in_slide(slide, pattern: str) -> Optional[str]:
    rx = re.compile(pattern)
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text or ""
        m = rx.search(txt)
        if m:
            return m.group(0)
    return None

def extract_server_name_from_slide(slide) -> Optional[str]:
    """
    템플릿에서 서버명이 들어가는 박스가 있다면:
    1) 이름 기반: SERVER_NAME
    2) fallback: 슬라이드 내 텍스트 중 'xxx-lnx' 같은 서버명 형태를 추정
    """
    shp = find_shape_by_name(slide, "SERVER_NAME")
    if shp and shp.has_text_frame:
        return _norm_server_name(shp.text_frame.text)

    # fallback: 슬라이드의 텍스트에서 서버명 후보를 찾음 (필요 시 패턴 조정)
    cand = find_text_in_slide(slide, r"[a-zA-Z0-9][a-zA-Z0-9\-_.]{2,}")
    return _norm_server_name(cand) if cand else None


# -----------------------------
# Table filling
# -----------------------------
def fill_stats_table(table, row_map: Dict[str, str], stats_by_metric: Dict[str, Stats], *,
                     pct_keys: set, bps_keys: set):
    """
    table: python-pptx Table
    row_map: { "표에 적힌 항목명": "metric_key" }
      예) {"CPU Used (%)": "cpu_used_pct", "MEM Used (%)": "mem_used_pct"}
    """
    # 컬럼 인덱스 가정: [0]=항목, [1]=최대, [2]=최소, [3]=평균  (템플릿에 맞게 조정)
    # 사진 예시와 같은 형태(항목/최대/최소/평균)라면 대체로 이 구조입니다.
    COL_ITEM, COL_MAX, COL_MIN, COL_AVG = 0, 1, 2, 3

    for r in range(1, len(table.rows)):  # 0행은 헤더라고 가정
        item = (table.cell(r, COL_ITEM).text or "").strip()
        if not item:
            continue

        metric_key = row_map.get(item)
        if not metric_key:
            # 템플릿의 항목명이 약간 다를 수 있어, 공백/기호 정규화 등을 원하면 여기서 처리
            continue

        st = stats_by_metric.get(metric_key, Stats())

        if metric_key in pct_keys:
            table.cell(r, COL_MAX).text = _fmt_pct(st.max)
            table.cell(r, COL_MIN).text = _fmt_pct(st.min)
            table.cell(r, COL_AVG).text = _fmt_pct(st.avg)
        elif metric_key in bps_keys:
            table.cell(r, COL_MAX).text = _fmt_bps(st.max)
            table.cell(r, COL_MIN).text = _fmt_bps(st.min)
            table.cell(r, COL_AVG).text = _fmt_bps(st.avg)
        else:
            table.cell(r, COL_MAX).text = _fmt(st.max, decimals=2)
            table.cell(r, COL_MIN).text = _fmt(st.min, decimals=2)
            table.cell(r, COL_AVG).text = _fmt(st.avg, decimals=2)


def find_all_tables(slide):
    tbls = []
    for shp in slide.shapes:
        if shp.has_table:
            tbls.append(shp.table)
    return tbls


# -----------------------------
# Image replacement
# -----------------------------
def replace_picture_in_bbox(slide, bbox_shape, image_path: str):
    """
    bbox_shape의 위치/크기에 맞춰 이미지 삽입.
    (기존 bbox_shape는 그대로 두거나 삭제할 수 있음. 여기서는 그대로 두고 위에 덮어씁니다.)
    """
    left, top, width, height = bbox_shape.left, bbox_shape.top, bbox_shape.width, bbox_shape.height
    pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    return pic


# -----------------------------
# Main population logic
# -----------------------------
def populate_ppt(
    template_pptx: str,
    output_pptx: str,
    reports: List[ServerReport],
    stats_by_server: Dict[str, Dict[str, Stats]],
    charts_by_server: Dict[str, Dict[str, str]],
):
    """
    charts_by_server 예:
      {
        "bastion-lnx": {
          "cpu_mem": "out/bastion-lnx_cpu_mem.png",
          "network": "out/bastion-lnx_network.png",
          "storage": "out/bastion-lnx_storage.png",
        }
      }

    템플릿에서 아래 shape name을 지정해두면 가장 확실합니다:
      SERVER_NAME
      CPU_MEM_CHART, NETWORK_CHART, STORAGE_CHART
      CPU_MEM_TABLE, NETWORK_TABLE, STORAGE_TABLE
    """
    prs = Presentation(template_pptx)

    # 템플릿 표의 "항목" 문구 → metric_key 매핑 (템플릿의 실제 텍스트와 정확히 맞추세요)
    CPU_MEM_ROW_MAP = {
        "CPU Used (%)": "cpu_used_pct",
        "MEM Used (%)": "mem_used_pct",
    }
    NETWORK_ROW_MAP = {
        "In bps (bps)": "net_in_bps",
        "Out bps (bps)": "net_out_bps",
    }
    STORAGE_ROW_MAP = {
        "파일시스템 사용률 (%)": "fs_used_pct",
    }

    pct_keys = {"cpu_used_pct", "mem_used_pct", "fs_used_pct"}
    bps_keys = {"net_in_bps", "net_out_bps"}

    for slide in prs.slides:
        server = extract_server_name_from_slide(slide)
        if not server:
            continue

        server_metrics = stats_by_server.get(server)
        server_charts = charts_by_server.get(server)

        if not server_metrics:
            # 서버명이 템플릿과 데이터에서 다르면 여기서 로그/예외 처리
            continue

        # ---- 1) 표 채우기 (이름 기반 우선) ----
        tbl_cpu = find_shape_by_name(slide, "CPU_MEM_TABLE")
        if tbl_cpu and tbl_cpu.has_table:
            fill_stats_table(tbl_cpu.table, CPU_MEM_ROW_MAP, server_metrics, pct_keys=pct_keys, bps_keys=bps_keys)

        tbl_net = find_shape_by_name(slide, "NETWORK_TABLE")
        if tbl_net and tbl_net.has_table:
            fill_stats_table(tbl_net.table, NETWORK_ROW_MAP, server_metrics, pct_keys=pct_keys, bps_keys=bps_keys)

        tbl_stg = find_shape_by_name(slide, "STORAGE_TABLE")
        if tbl_stg and tbl_stg.has_table:
            fill_stats_table(tbl_stg.table, STORAGE_ROW_MAP, server_metrics, pct_keys=pct_keys, bps_keys=bps_keys)

        # ---- 2) 그래프 이미지 삽입 (이름 기반 우선) ----
        if server_charts:
            shp_cpu = find_shape_by_name(slide, "CPU_MEM_CHART")
            if shp_cpu and server_charts.get("cpu_mem"):
                replace_picture_in_bbox(slide, shp_cpu, server_charts["cpu_mem"])

            shp_net = find_shape_by_name(slide, "NETWORK_CHART")
            if shp_net and server_charts.get("network"):
                replace_picture_in_bbox(slide, shp_net, server_charts["network"])

            shp_stg = find_shape_by_name(slide, "STORAGE_CHART")
            if shp_stg and server_charts.get("storage"):
                replace_picture_in_bbox(slide, shp_stg, server_charts["storage"])

        # ---- 3) (옵션) 서버명 텍스트 표준화 ----
        shp_name = find_shape_by_name(slide, "SERVER_NAME")
        if shp_name and shp_name.has_text_frame:
            shp_name.text_frame.text = server

    prs.save(output_pptx)


# -----------------------------
# Example usage
# -----------------------------
if __name__ == "__main__":
    # 예시 데이터
    stats_by_server = {
        "bastion-lnx": {
            "cpu_used_pct": Stats(min=0.39, max=0.92, avg=0.48),
            "mem_used_pct": Stats(min=12.76, max=14.68, avg=13.01),
            "net_in_bps": Stats(min=241.96, max=423170.0, avg=16540.0),
            "net_out_bps": Stats(min=3010.0, max=133410.0, avg=16240.0),
            "fs_used_pct": Stats(min=7.68, max=7.78, avg=7.73),
        }
    }

    charts_by_server = {
        "bastion-lnx": {
            "cpu_mem": "charts/bastion-lnx_cpu_mem.png",
            "network": "charts/bastion-lnx_network.png",
            "storage": "charts/bastion-lnx_storage.png",
        }
    }

    populate_ppt(
        template_pptx="template.pptx",
        output_pptx="report.pptx",
        stats_by_server=stats_by_server,
        charts_by_server=charts_by_server,
    )