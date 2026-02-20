from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

def change_shape_text(
    pptx_path: str,
    shape_name: str,
    new_text: str,
    font_size: int,
    font_name: str = "맑은 고딕",
    font_align: str = "left",
    font_bold: bool = False,
    font_color_rgb: RGBColor = RGBColor(0x00, 0x00, 0x00),
    slide_index: int = 0
):
    """
    pptx_path   : 수정할 pptx 파일
    shape_name  : 찾을 shape.name (예: '기관명')
    new_text    : 교체할 문자열
    slide_index : 기본 0 (첫 번째 표지)
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    target = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target = shape
            break

    if target is None:
        raise ValueError(f"'{shape_name}' 이름을 가진 shape을 찾을 수 없습니다.")

    if not target.has_text_frame:
        raise ValueError(f"'{shape_name}' shape에 텍스트 프레임이 없습니다.")

    tf = target.text_frame
    tf.clear()
    tf.text = new_text
    
    # 스타일 설정
    p = tf.paragraphs[0]
    if font_align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif font_align == "left":
        p.alignment = PP_ALIGN.LEFT
    elif font_align == "right":
        p.alignment = PP_ALIGN.RIGHT
    
    run = p.runs[0]
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.color.rgb = font_color_rgb
    font.bold = font_bold
    prs.save(pptx_path)


def copy_pptx_to_multiple_names(
    date: str, 
    src_pptx: str | Path,
    target_pptx: str | Path,
    output_dir: str | Path = None,
):
    """
    src_pptx: 원본 pptx 파일 (a)
    target_pptx: 복사본 파일 경로 (b)
    """

    src_pptx = Path(src_pptx)
    output_dir = Path(output_dir) if output_dir else Path.cwd()

    filename = f"{src_pptx.stem}({target_pptx}){src_pptx.suffix}"
    file_path = output_dir / filename

    # 1) 파일 복사
    file_path.write_bytes(src_pptx.read_bytes())

    # 2) 표지의 '기관명' 텍스트 교체
    change_shape_text(
        pptx_path=str(file_path),
        shape_name="기관명",
        new_text=target_pptx,
        font_size=36,
        font_align="left",
        font_bold=False,
        font_color_rgb=RGBColor(0x00, 0x42, 0xFF)
    )

    #3) 날짜 텍스트 교체
    change_shape_text(
        pptx_path=str(file_path),
        shape_name="날짜",
        new_text=date,
        font_size=12,
        font_align="left",
        font_bold=False,
        font_color_rgb=RGBColor(0x00, 0x21, 0x46)
    )

    print(f"기관별 ppt 생성 완료: {file_path}")