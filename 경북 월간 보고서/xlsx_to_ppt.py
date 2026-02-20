from __future__ import annotations

from dataclasses import dataclass
import io
from typing import Dict, Optional, List

from pptx import Presentation


from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER

from copy import deepcopy

# -----------------------------
# Data model
# -----------------------------
@dataclass
class Stats:
    min: Optional[float] = None
    max: Optional[float] = None
    avg: Optional[float] = None

@dataclass
class MetricBlock:
    # CPU/MEM 은 2개 지표, Network는 IN/OUT 2개 지표 등 복수 가능
    stats_by_metric: Dict[str, Stats]     # e.g. {"CPU": Stats(...), "MEM": Stats(...)}
    images: List[bytes]                  # 이미지(PNG/JPG 바이너리). 블록 내 이미지 순서대로 담김

@dataclass
class ServerReport:
    server_name: str
    cpu_mem: MetricBlock
    network: MetricBlock
    filesystem: MetricBlock

# -----------------------------
# PPT 이미지 삽입 함수
# -----------------------------
def get_table_cell_bbox(table_shape, row: int, col_start: int, col_end: int = None):
    """
    table_shape: slide.shapes 중 has_table == True 인 shape (table을 포함한 shape)
    row, col: 0-based index
    return: (left, top, width, height) 모두 EMU 단위
    """
    tbl = table_shape.table

    # 테이블 시작 좌표
    left = table_shape.left
    top = table_shape.top

    # col 이전 열들의 너비를 누적
    for c in range(col_start):
        left += tbl.columns[c].width

    # top은 수동으로...
    if row == 1:
        top = Cm(3.57)    
    elif row == 6:
        top = Cm(8.49)
    elif row == 11:
        top = Cm(13.41)

    if col_end is not None:
        # col_start ~ col_end-1 열들의 너비를 누적
        width = Emu(0)
        for c in range(col_start, col_end+1):
            width += tbl.columns[c].width
    height = tbl.rows[row].height

    return left, top, width, height

def add_picture_over_table_cell(slide, table_shape, image_bytes: bytes, row: int, col_start: int, col_end: int = None):
    bio = io.BytesIO(image_bytes)

    left, top, width, height = get_table_cell_bbox(table_shape, row, col_start, col_end)
    # left, top, width, height = get_table_cell_bbox_scaled(table_shape, row, col_start, col_end)

    pic = slide.shapes.add_picture(
        bio,
        left,
        top,
        width=width,
        height=height
    )
    return pic


# -----------------------------
# PPT 테이블 찾기 함수
# -----------------------------
def find_main_table(slide):
    """
    테이블 찾는 함수
    
    :param slide: 슬라이드 객체(페이지)
    """
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break
    return table_shape

# ---------------------
# Main population logic
# -----------------------------
def set_cell_text_style(
        cell, 
        align, 
        font_name, 
        font_size, 
        text: str,
        is_bold: bool = False,
        font_color: RGBColor = RGBColor(0x00, 0x00, 0x00)
):
    tf = cell.text_frame
    tf.clear()  # 기존 문단/Run 제거

    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER  # 가운데 정렬
    elif align == "left":
        p.alignment = PP_ALIGN.LEFT    # 좌측 정렬
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT  # 우측 정렬

    run = p.add_run()
    run.text = text
    
    font = run.font
    # font.name = "맑은 고딕"
    # font.size = Pt(8)
    font.name = font_name
    font.size = Pt(font_size)
    font.color.rgb = font_color
    if is_bold:
        font.bold = True


def populate_slide_with_report(
    slide,
    report: ServerReport
):
    # ---- 1) 표 찾기 ----
    table_shape = find_main_table(slide)
    tbl = table_shape.table

    # ---- 2) 서버명 삽입 ----
    set_cell_text_style(tbl.cell(0,0), "center", "맑은 고딕", 10, report.server_name, True, RGBColor(0xFF, 0xFF, 0xFF))

    # ---- 2) 표 채우기 및 그래프 이미지 삽입 (이름 기반 우선) ----
    cpu_mem = report.cpu_mem
    cpu_stats = cpu_mem.stats_by_metric["● CPU Used (%)"]
    mem_stats = cpu_mem.stats_by_metric["● MEM Used (%)"]
    cpu_mem_images = cpu_mem.images
    set_cell_text_style(tbl.cell(3, 2), "right", "맑은 고딕", 8, cpu_stats.max)
    set_cell_text_style(tbl.cell(3, 3), "right", "맑은 고딕", 8, cpu_stats.min)
    set_cell_text_style(tbl.cell(3, 4), "right", "맑은 고딕", 8, cpu_stats.avg)
    set_cell_text_style(tbl.cell(4, 2), "right", "맑은 고딕", 8, mem_stats.max)
    set_cell_text_style(tbl.cell(4, 3), "right", "맑은 고딕", 8, mem_stats.min)
    set_cell_text_style(tbl.cell(4, 4), "right", "맑은 고딕", 8, mem_stats.avg)
    add_picture_over_table_cell(slide, table_shape, image_bytes=cpu_mem_images[0], row=1, col_start=1, col_end=4)

    network = report.network
    network_stats = network.stats_by_metric
    network_images = network.images
    set_cell_text_style(tbl.cell(8, 2), "right", "맑은 고딕", 8, network_stats["● In bps (bps)"].max)
    set_cell_text_style(tbl.cell(8, 3), "right", "맑은 고딕", 8, network_stats["● In bps (bps)"].min)
    set_cell_text_style(tbl.cell(8, 4), "right", "맑은 고딕", 8, network_stats["● In bps (bps)"].avg)
    set_cell_text_style(tbl.cell(9, 2), "right", "맑은 고딕", 8, network_stats["● Out bps (bps)"].max)
    set_cell_text_style(tbl.cell(9, 3), "right", "맑은 고딕", 8, network_stats["● Out bps (bps)"].min)
    set_cell_text_style(tbl.cell(9, 4), "right", "맑은 고딕", 8, network_stats["● Out bps (bps)"].avg)
    add_picture_over_table_cell(slide, table_shape, image_bytes=network_images[0], row=6, col_start=1, col_end=4)
    # replace_picture_in_bbox(tbl.cell(6, 1), network_images[0])

    filesystem = report.filesystem
    filesystem_stats = filesystem.stats_by_metric
    filesystem_images = filesystem.images
    set_cell_text_style(tbl.cell(13, 2), "right", "맑은 고딕", 8, filesystem_stats["● 파일시스템 사용률 (%)"].max)
    set_cell_text_style(tbl.cell(13, 3), "right", "맑은 고딕", 8, filesystem_stats["● 파일시스템 사용률 (%)"].min)
    set_cell_text_style(tbl.cell(13, 4), "right", "맑은 고딕", 8, filesystem_stats["● 파일시스템 사용률 (%)"].avg)
    add_picture_over_table_cell(slide, table_shape, image_bytes=filesystem_images[0], row=11, col_start=1, col_end=4)
    # replace_picture_in_bbox(tbl.cell(11, 1), report.filesystem.images[0])

# -----------------------------
# PPT 제목 플레이스홀더 클리어 함수
# -----------------------------
def clear_title_placeholders(slide):
    for shp in slide.shapes:
        if not shp.is_placeholder:
            continue
        ph = shp.placeholder_format
        if ph.type == PP_PLACEHOLDER.TITLE:
            if shp.has_text_frame:
                shp.text_frame.clear()

# -----------------------------
# PPT 슬라이드 복제 함수
# -----------------------------
def duplicate_slide(prs, slide_index=0):
    """
    prs: Presentation
    slide_index: 복제할 템플릿 슬라이드 인덱스 (보통 0)
    return: 새로 생성된 슬라이드 객체
    """
    source = prs.slides[slide_index]
    layout = source.slide_layout  # 원본과 동일 레이아웃 유지
    new_slide = prs.slides.add_slide(layout)

    # 레이아웃이 만든 기본 제목 제거
    clear_title_placeholders(new_slide)

    # 원본 슬라이드의 모든 shape를 새 슬라이드로 복사
    for shp in source.shapes:
        new_el = deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


# -----------------------------
# PPT 빌드 함수 (안전 버전)
# -----------------------------
def build_ppt_from_template(template_pptx, output_pptx,reports, date):

    prs = Presentation(template_pptx)
    if not reports:
        raise ValueError("reports가 비어 있습니다.")
    
    # 날짜 수정
    for shape in prs.slides[0].shapes:
        if not shape.name=="날짜": continue
        if not shape.has_text_frame:continue
        tf = shape.text_frame
        p = tf.paragraphs[0]
        # 런이 없으면 하나 만들기
        if not p.runs:
            run = p.add_run()
        else:
            run = p.runs[0]
        # 나머지 런들은 빈 문자열로 만들어 사실상 "삭제" 효과
        for r in p.runs[1:]:
            r.text = ""
        run.text = f"(기간: {date})"
        

    # 템플릿 슬라이드(0번)를 report 수만큼 확장
    while len(prs.slides) < len(reports):
        duplicate_slide(prs, slide_index=0)

    # 각 슬라이드에 각 report 매핑
    for i, rep in enumerate(reports):
        populate_slide_with_report(prs.slides[i], rep)

    prs.save(output_pptx)
    print(f"PPTX 파일 생성 완료: {output_pptx}")


# -----------------------------
# Example usage
# -----------------------------
# if __name__ == "__main__":
