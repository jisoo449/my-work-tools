from pptx import Presentation

prs = Presentation("template02.pptx")

for slide_idx, slide in enumerate(prs.slides, start=1):
    table = None

    # 1. 슬라이드 내 테이블 탐색
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print(f"Slide {slide_idx}: 테이블 없음")
        continue

    # 2. 행/열 개수 확인
    row_count = len(table.rows)
    col_count = len(table.columns)

    print(f"Slide {slide_idx}: 테이블 크기 = {row_count}행 x {col_count}열")
